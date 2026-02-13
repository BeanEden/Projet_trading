
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Projet Fil Rouge : Système de Trading GBP/USD"
    subtitle.text = "M1 → M15 → ML → RL → API → Docker\nFévrier 2026"

    # 2. Contexte & Objectifs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Contexte & Objectifs"
    content = slide.placeholders[1].text_frame
    content.text = "Conception d'un système de décision algorithmique sur la paire GBP/USD."
    
    p = content.add_paragraph()
    p.text = "Données : Fréquence brute 1 minute -> Décision 15 minutes."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Actions : BUY, SELL, HOLD."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Objectif : Maximiser le profit cumulé sous contraintes :"
    p.level = 1
    
    sub = content.add_paragraph()
    sub.text = "- Coûts de transaction réalistes"
    sub.level = 2
    sub = content.add_paragraph()
    sub.text = "- Drawdown limité"
    sub.level = 2
    sub = content.add_paragraph()
    sub.text = "- Robustesse inter-annuelle"
    sub.level = 2

    # 3. Données & Protocole
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Données & Split Temporel"
    content = slide.placeholders[1].text_frame
    content.text = "Période disponible : 2022 - 2024"
    
    p = content.add_paragraph()
    p.text = "Protocole de Split Strict (Interdiction de split aléatoire) :"
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "2022 : Entraînement (Train)"
    p.level = 2
    p = content.add_paragraph()
    p.text = "2023 : Validation / Optimisation"
    p.level = 2
    p = content.add_paragraph()
    p.text = "2024 : Test Final (Jamais vu par le modèle)"
    p.level = 2

    # 4. Pipeline Data Science
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Pipeline Data Science"
    content = slide.placeholders[1].text_frame
    
    p = content.add_paragraph()
    p.text = "Phase 1 : Importation M1"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Nettoyage, fusion date/heure, vérification régularité."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Phase 2 : Agrégation M1 -> M15"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Resampling (Open=First, High=Max, Low=Min, Close=Last)."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Phase 3 : Nettoyage"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Filtre des bougies incomplètes et gaps anormaux."
    p.level = 1

    # 5. Feature Engineering V2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Feature Engineering V2"
    content = slide.placeholders[1].text_frame
    content.text = "Indicateurs techniques calculés sur le passé uniquement."
    
    p = content.add_paragraph()
    p.text = "Bloc Court Terme (Dynamique) :"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Returns, EMA (20, 50), RSI 14, Volatilité roulante."
    p.level = 2
    
    p = content.add_paragraph()
    p.text = "Bloc Contexte (Régime) :"
    p.level = 1
    p = content.add_paragraph()
    p.text = "Tendance long terme (EMA 200), ATR 14, ADX, MACD."
    p.level = 2

    # 6. Stratégies & Modélisation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Stratégies & Modélisation"
    content = slide.placeholders[1].text_frame
    
    p = content.add_paragraph()
    p.text = "1. Baseline : Règles Fixes"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Stratégie simple (ex: EMA Cross + RSI) pour établir un niveau de référence."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "2. Machine Learning"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Objectif : Prédire la direction (Hausse/Baisse) de la prochaine bougie."
    p.level = 1
    p = content.add_paragraph()
    p.text = "Algorithmes : Random Forest, Gradient Boosting."
    p.level = 1

    # 7. Architecture Technique
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Architecture Technique & Industrialisation"
    content = slide.placeholders[1].text_frame
    
    p = content.add_paragraph()
    p.text = "Backend : API FastAPI"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Expose le meilleur modèle, gère les inférences."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "Frontend : Web App Flask"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Dashboard utilisateur, visualisation des prédictions."
    p.level = 1
    
    p = content.add_paragraph()
    p.text = "MLOps :"
    p.level = 0
    p = content.add_paragraph()
    p.text = "Model Registry (Versioning v1/v2), Dockerisation."
    p.level = 1

    # 8. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    content = slide.placeholders[1].text_frame
    content.text = "Un modèle performant n'est pas celui qui gagne le plus sur 2022, mais celui qui :"
    
    p = content.add_paragraph()
    p.text = "Survit au changement de régime (Robustesse 2024)."
    p.level = 1
    p = content.add_paragraph()
    p.text = "Tient compte des coûts de transaction."
    p.level = 1
    p = content.add_paragraph()
    p.text = "Évite l'overfitting temporel."
    p.level = 1
    p = content.add_paragraph()
    p.text = "Est reproductible et industrialisable."
    p.level = 1

    # Save
    output_path = "Presentation_Projet_Trading.pptx"
    prs.save(output_path)
    print(f"Présentation générée : {output_path}")

if __name__ == "__main__":
    create_presentation()
